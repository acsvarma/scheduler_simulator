"""
SIRIUS CAR-T Scheduler — Presentation Generator
Run:  python generate_presentation.py
Output: SIRIUS_Scheduler_Overview.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
import copy
from lxml import etree

# ── Brand palette ──────────────────────────────────────────────────────────────
BMS_PURPLE  = RGBColor(0xBE, 0x2B, 0xBB)
BMS_DARK    = RGBColor(0x1A, 0x1A, 0x2E)
BMS_MID     = RGBColor(0x2C, 0x3E, 0x50)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xF4, 0xF6, 0xF8)
MID_GRAY    = RGBColor(0x95, 0xA5, 0xA6)
DARK_GRAY   = RGBColor(0x34, 0x39, 0x44)
BLUE        = RGBColor(0x29, 0x80, 0xB9)
GREEN       = RGBColor(0x27, 0xAE, 0x60)
ORANGE      = RGBColor(0xE6, 0x7E, 0x22)
RED         = RGBColor(0xC0, 0x39, 0x2B)
TEAL        = RGBColor(0x16, 0xA0, 0x85)
GOLD        = RGBColor(0xF3, 0x9C, 0x12)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Low-level helpers ──────────────────────────────────────────────────────────

def rgb(r, g, b):
    return RGBColor(r, g, b)


def add_rect(slide, left, top, width, height, fill_color, border_color=None, border_pt=0):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color and border_pt:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_name="Trebuchet MS", font_size=12, bold=False, italic=False,
                 color=DARK_GRAY, align=PP_ALIGN.LEFT, word_wrap=True, v_anchor="top"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_para(tf, text, font_size=11, bold=False, color=DARK_GRAY,
             font_name="Trebuchet MS", align=PP_ALIGN.LEFT, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p


def add_bullet(tf, text, level=0, font_size=11, color=DARK_GRAY, bold=False):
    p = tf.add_paragraph()
    p.level = level
    run = p.add_run()
    run.text = text
    run.font.name = "Trebuchet MS"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def header_bar(slide, title, subtitle=None):
    """Purple top bar with title."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.15), BMS_PURPLE)
    add_text_box(slide, title,
                 Inches(0.4), Inches(0.12), Inches(11), Inches(0.6),
                 font_size=26, bold=True, color=WHITE)
    if subtitle:
        add_text_box(slide, subtitle,
                     Inches(0.4), Inches(0.72), Inches(11), Inches(0.4),
                     font_size=13, color=WHITE)
    # BMS badge top-right
    add_text_box(slide, "BMS  |  SIRIUS",
                 Inches(11.2), Inches(0.3), Inches(2), Inches(0.5),
                 font_size=11, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)


def footer_bar(slide, page_num, total):
    add_rect(slide, 0, Inches(7.1), SLIDE_W, Inches(0.4), BMS_DARK)
    add_text_box(slide, "CONFIDENTIAL — Bristol-Myers Squibb  |  SIRIUS CAR-T Scheduler",
                 Inches(0.3), Inches(7.12), Inches(10), Inches(0.3),
                 font_size=8, color=MID_GRAY)
    add_text_box(slide, f"{page_num} / {total}",
                 Inches(12.5), Inches(7.12), Inches(0.7), Inches(0.3),
                 font_size=8, color=MID_GRAY, align=PP_ALIGN.RIGHT)


def card(slide, left, top, width, height, title, bullets,
         title_color=BMS_PURPLE, bg=LIGHT_GRAY, border=None):
    add_rect(slide, left, top, width, height, bg,
             border_color=border or title_color, border_pt=1.5)
    add_text_box(slide, title,
                 left + Inches(0.12), top + Inches(0.08),
                 width - Inches(0.2), Inches(0.38),
                 font_size=11, bold=True, color=title_color)
    txBox = slide.shapes.add_textbox(
        left + Inches(0.12), top + Inches(0.46),
        width - Inches(0.2), height - Inches(0.56))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = b
        run.font.name = "Trebuchet MS"
        run.font.size = Pt(9.5)
        run.font.color.rgb = DARK_GRAY


def pill(slide, left, top, width, height, text, bg, text_color=WHITE, font_size=10):
    shape = slide.shapes.add_shape(5, left, top, width, height)  # ROUNDED_RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Trebuchet MS"
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = text_color


# ── Slide builders ─────────────────────────────────────────────────────────────

def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Full purple background
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, BMS_PURPLE)
    # Dark bottom strip
    add_rect(slide, 0, Inches(5.8), SLIDE_W, Inches(1.7), BMS_DARK)

    # BMS / SIRIUS badge
    add_rect(slide, Inches(0.5), Inches(0.4), Inches(2.2), Inches(0.55),
             WHITE, border_color=WHITE, border_pt=0)
    add_text_box(slide, "BMS  ·  Bristol-Myers Squibb",
                 Inches(0.55), Inches(0.42), Inches(2.1), Inches(0.5),
                 font_size=9, bold=True, color=BMS_PURPLE)

    # Main title
    add_text_box(slide, "SIRIUS",
                 Inches(0.7), Inches(1.3), Inches(12), Inches(1.3),
                 font_size=72, bold=True, color=WHITE)
    add_text_box(slide, "CAR-T Production Scheduler",
                 Inches(0.7), Inches(2.5), Inches(12), Inches(0.7),
                 font_size=30, bold=False, color=WHITE)

    # Separator line
    add_rect(slide, Inches(0.7), Inches(3.35), Inches(4), Inches(0.04), WHITE)

    add_text_box(slide,
                 "Scheduling, Simulation & Resource Optimisation\n"
                 "for Autologous Cell Therapy Manufacturing",
                 Inches(0.7), Inches(3.5), Inches(11), Inches(0.9),
                 font_size=14, color=WHITE)

    # Bottom strip text
    add_text_box(slide, "Cell Therapy Operations  ·  Manufacturing Sciences  ·  Confidential",
                 Inches(0.7), Inches(6.05), Inches(10), Inches(0.4),
                 font_size=11, color=MID_GRAY)
    add_text_box(slide, "2026",
                 Inches(11.8), Inches(6.05), Inches(1.2), Inches(0.4),
                 font_size=11, color=MID_GRAY, align=PP_ALIGN.RIGHT)

    return slide


def slide_agenda(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, "Agenda", "What this presentation covers")
    footer_bar(slide, page, total)

    topics = [
        ("01", "CAR-T Manufacturing Overview",    "What CAR-T is and why scheduling is complex"),
        ("02", "The 10-Step Process",              "Upstream steps from TSA through Fill"),
        ("03", "Equipment & Robots",               "Three ROMAG, 47 incubators, three robot resources"),
        ("04", "Scheduling Rules & Constraints",   "Biological minimums, chained steps, incubator locking"),
        ("05", "Incubator Pool Management",        "47-slot queuing, continuous reservation model"),
        ("06", "Equipment Efficiency",             "OEE factor and its effect on schedule & throughput"),
        ("07", "Throughput & Bottleneck Analysis", "Theoretical max, AMS throughput, robot contention"),
        ("08", "Controls & Monitoring",            "Simulation engine, live mode, SAP/DeltaV interfaces"),
    ]

    col_w = Inches(5.8)
    for i, (num, title, sub) in enumerate(topics):
        row, col = divmod(i, 2)
        left = Inches(0.35) + col * Inches(6.7)
        top  = Inches(1.35) + row * Inches(1.42)

        add_rect(slide, left, top, col_w, Inches(1.25), LIGHT_GRAY,
                 border_color=BMS_PURPLE, border_pt=1.5)
        add_text_box(slide, num,
                     left + Inches(0.12), top + Inches(0.1),
                     Inches(0.55), Inches(0.45),
                     font_size=22, bold=True, color=BMS_PURPLE)
        add_text_box(slide, title,
                     left + Inches(0.75), top + Inches(0.1),
                     col_w - Inches(0.85), Inches(0.45),
                     font_size=13, bold=True, color=BMS_DARK)
        add_text_box(slide, sub,
                     left + Inches(0.75), top + Inches(0.55),
                     col_w - Inches(0.85), Inches(0.55),
                     font_size=10, color=DARK_GRAY)
    return slide


def slide_cart_overview(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, "CAR-T Manufacturing Overview",
               "Why scheduling autologous cell therapy is uniquely complex")
    footer_bar(slide, page, total)

    # Left column — what is CAR-T
    add_rect(slide, Inches(0.35), Inches(1.3), Inches(5.9), Inches(5.7), LIGHT_GRAY,
             border_color=BMS_PURPLE, border_pt=1)
    add_text_box(slide, "What is CAR-T?",
                 Inches(0.5), Inches(1.4), Inches(5.6), Inches(0.45),
                 font_size=14, bold=True, color=BMS_PURPLE)

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(5.6), Inches(4.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    items = [
        "Chimeric Antigen Receptor T-cell (CAR-T) therapy is a personalised cancer treatment "
        "in which a patient's own T-cells are extracted, genetically re-engineered in a "
        "manufacturing facility, and re-infused to target and destroy cancer cells.",
        "",
        "Each batch is PATIENT-SPECIFIC:",
        "  •  Cells collected from one patient",
        "  •  Processed through a fixed sequence of steps",
        "  •  Re-administered to that same patient",
        "",
        "Every patient's material must be tracked independently from vein-to-vein. "
        "No pooling, no substitution.",
    ]
    for i, txt in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = txt
        run.font.name = "Trebuchet MS"
        run.font.size = Pt(10.5 if not txt.startswith("Each") and not txt.startswith("Every") else 10.5)
        run.font.bold = txt.endswith(":")
        run.font.color.rgb = DARK_GRAY

    # Right column — scheduling challenges
    add_rect(slide, Inches(6.55), Inches(1.3), Inches(6.43), Inches(5.7), BMS_DARK,
             border_color=BMS_DARK, border_pt=0)
    add_text_box(slide, "Scheduling Challenges",
                 Inches(6.7), Inches(1.4), Inches(6.1), Inches(0.45),
                 font_size=14, bold=True, color=BMS_PURPLE)

    challenges = [
        ("Long cycle time",       "~7–10 days per patient from leukapheresis to fill"),
        ("Fixed sequence",        "10 steps must occur in strict order per patient"),
        ("Shared equipment",      "3 ROMAGs, 47 incubators, 3 robots — all shared across patients"),
        ("Robot serialisation",   "One robot handles one operation at a time; queuing is critical"),
        ("Incubator dedication",  "Each patient holds their incubator for the full upstream period"),
        ("DHWF → FILL chain",     "Harvest must flow directly to Fill — no intermediate hold allowed"),
        ("Biological floors",     "Minimum 2 h incubation before patient can advance"),
        ("47-patient cap",        "Limited by incubator count; excess patients queue until slot opens"),
    ]
    for i, (title, desc) in enumerate(challenges):
        top = Inches(1.95) + i * Inches(0.61)
        add_rect(slide, Inches(6.7), top, Inches(5.9), Inches(0.52),
                 rgb(0x23, 0x2B, 0x3B), border_color=BMS_PURPLE, border_pt=0.5)
        add_text_box(slide, title,
                     Inches(6.82), top + Inches(0.04),
                     Inches(1.6), Inches(0.45),
                     font_size=9.5, bold=True, color=BMS_PURPLE)
        add_text_box(slide, desc,
                     Inches(8.45), top + Inches(0.04),
                     Inches(4.1), Inches(0.45),
                     font_size=9.5, color=WHITE)

    return slide


def slide_process(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, "The 10-Step Upstream Process",
               "Fixed sequence from T-cell activation to product fill — every patient follows this path")
    footer_bar(slide, page, total)

    steps = [
        ("UP_TSA",   "T-Cell\nActivation",   "4 h",   BLUE,   "ROMAG",     "Romag Robot"),
        ("UP_INC 1", "Incubation\n(48 h)",   "48 h",  TEAL,   "Incubator", "—"),
        ("UP_BSD 1", "Bead Selection\n1st",  "2 h",   ORANGE, "BSD",       "—"),
        ("UP_INC 2", "Incubation\n(24 h)",   "24 h",  TEAL,   "Incubator", "—"),
        ("UP_TxD",   "Transduction",         "6 h",   BMS_PURPLE, "ROMAG", "Romag Robot"),
        ("UP_INC 3", "Incubation\n(48 h)",   "48 h",  TEAL,   "Incubator", "—"),
        ("UP_BSD 2", "Bead Selection\n2nd",  "2 h",   ORANGE, "BSD",       "—"),
        ("UP_INC 4", "Incubation\n(24 h)",   "24 h",  TEAL,   "Incubator", "—"),
        ("UP_DHWF",  "Harvest &\nFormulate", "8 h",   GREEN,  "ROMAG",     "Romag + Inc Robot"),
        ("UP_FILL",  "Fill &\nFinish",       "3 h",   RED,    "Fill Stn",  "Fill Robot"),
    ]

    box_w = Inches(1.18)
    box_h = Inches(1.55)
    gap   = Inches(0.1)
    top   = Inches(1.45)
    start_left = Inches(0.3)

    for i, (code, name, dur, color, eq, robot) in enumerate(steps):
        left = start_left + i * (box_w + gap)

        # Arrow connector (except last)
        if i < len(steps) - 1:
            arr_left = left + box_w
            add_rect(slide, arr_left, top + Inches(0.65), gap, Inches(0.1), MID_GRAY)

        # Step box
        add_rect(slide, left, top, box_w, box_h, color)

        # Step code
        add_text_box(slide, code,
                     left + Inches(0.05), top + Inches(0.05),
                     box_w - Inches(0.1), Inches(0.32),
                     font_size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Step name
        add_text_box(slide, name,
                     left + Inches(0.05), top + Inches(0.35),
                     box_w - Inches(0.1), Inches(0.62),
                     font_size=9.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Duration
        add_rect(slide, left + Inches(0.1), top + Inches(1.02),
                 box_w - Inches(0.2), Inches(0.28), rgb(0, 0, 0))
        add_text_box(slide, dur,
                     left + Inches(0.1), top + Inches(1.03),
                     box_w - Inches(0.2), Inches(0.28),
                     font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Equipment label below box
        add_text_box(slide, eq,
                     left, top + box_h + Inches(0.08),
                     box_w, Inches(0.28),
                     font_size=8.5, bold=False, color=DARK_GRAY, align=PP_ALIGN.CENTER)
        # Robot label
        add_text_box(slide, robot,
                     left, top + box_h + Inches(0.35),
                     box_w, Inches(0.28),
                     font_size=7.5, color=BMS_PURPLE, align=PP_ALIGN.CENTER)

    # Legend
    legend_items = [
        (BLUE,       "ROMAG step"),
        (TEAL,       "Incubation"),
        (ORANGE,     "BSD Sampling"),
        (GREEN,      "Harvest"),
        (RED,        "Fill"),
    ]
    add_text_box(slide, "Legend:",
                 Inches(0.35), Inches(6.3), Inches(1.0), Inches(0.3),
                 font_size=9, bold=True, color=DARK_GRAY)
    for i, (c, lbl) in enumerate(legend_items):
        lx = Inches(1.35) + i * Inches(2.35)
        add_rect(slide, lx, Inches(6.38), Inches(0.22), Inches(0.18), c)
        add_text_box(slide, lbl, lx + Inches(0.28), Inches(6.33),
                     Inches(2.0), Inches(0.28), font_size=8.5, color=DARK_GRAY)

    # Key constraint callouts
    add_rect(slide, Inches(0.3), Inches(5.65), Inches(12.7), Inches(0.52),
             rgb(0xFD, 0xF2, 0xE9), border_color=ORANGE, border_pt=1)
    add_text_box(slide,
                 "⚠  Key Constraints:   "
                 "INC steps → 2 h biological minimum before advancement   |   "
                 "Incubator held from UP_TSA start → UP_DHWF start (entire upstream)   |   "
                 "DHWF must chain directly to FILL (no hold)",
                 Inches(0.45), Inches(5.67), Inches(12.4), Inches(0.45),
                 font_size=9, color=rgb(0x8B, 0x4A, 0x00))

    return slide


def slide_equipment(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, "Equipment & Robot Resources",
               "Shared assets across all concurrent patient batches")
    footer_bar(slide, page, total)

    # Equipment table
    add_text_box(slide, "Equipment Pool",
                 Inches(0.35), Inches(1.3), Inches(7.5), Inches(0.4),
                 font_size=14, bold=True, color=BMS_DARK)

    eq_data = [
        ("Equipment Type",  "Count",  "Used By",             "Shared?"),
        ("ROMAG",           "3",      "UP_TSA, UP_TxD, UP_DHWF", "Yes — up to 3 patients simultaneously"),
        ("Incubator",       "47",     "All 4 INC steps",     "No — 1 dedicated per patient"),
        ("BSD Sampling",    "1",      "UP_BSD_1, UP_BSD_2",  "Yes"),
        ("BSD Weighing",    "1",      "Support steps",       "Yes"),
        ("Fill Station",    "1",      "UP_FILL",             "Yes — single serialised resource"),
    ]

    col_widths = [Inches(1.8), Inches(0.7), Inches(2.6), Inches(2.9)]
    col_x = [Inches(0.35), Inches(2.2), Inches(2.95), Inches(5.6)]
    row_h = Inches(0.52)

    for r, row in enumerate(eq_data):
        bg = BMS_DARK if r == 0 else (LIGHT_GRAY if r % 2 == 0 else WHITE)
        txt_c = WHITE if r == 0 else DARK_GRAY
        bold = r == 0
        for c, (val, cw, cx) in enumerate(zip(row, col_widths, col_x)):
            add_rect(slide, cx, Inches(1.75) + r * row_h, cw, row_h,
                     bg, border_color=MID_GRAY, border_pt=0.5)
            add_text_box(slide, val, cx + Inches(0.08),
                         Inches(1.75) + r * row_h + Inches(0.1),
                         cw - Inches(0.12), row_h - Inches(0.1),
                         font_size=10, bold=bold, color=txt_c)

    # Robot resources section
    add_text_box(slide, "Three Independent Robot Resources",
                 Inches(0.35), Inches(5.08), Inches(12.5), Inches(0.4),
                 font_size=14, bold=True, color=BMS_DARK)

    robots = [
        ("🦾 Romag Robot",
         BLUE,
         ["Used by: UP_TSA, UP_TxD, UP_DHWF",
          "Segments: Cell Labelling, Transduction Prep,",
          "  Harvest Prep",
          "Robot time: 5–10 min per segment",
          "Serialised: only one patient at a time"]),
        ("🔬 Incubation Robot",
         TEAL,
         ["Used by: Material transfer segments",
          "Appears in: UP_TSA Seg 3 & 5,",
          "  UP_DHWF Seg 5",
          "Robot time: 3–10 min per transfer",
          "Serialised: only one patient at a time"]),
        ("💊 Fill Robot",
         RED,
         ["Used by: UP_FILL Filling segment",
          "Robot time: 120 min (full fill duration)",
          "Longest single robot occupation",
          "Serialised: only one patient at a time",
          "Constraint: Fill Station + Fill Robot both needed"]),
    ]

    for i, (name, color, bullets) in enumerate(robots):
        left = Inches(0.35) + i * Inches(4.3)
        add_rect(slide, left, Inches(5.55), Inches(4.1), Inches(1.6),
                 LIGHT_GRAY, border_color=color, border_pt=2)
        add_text_box(slide, name, left + Inches(0.12), Inches(5.62),
                     Inches(3.8), Inches(0.38),
                     font_size=12, bold=True, color=color)
        txBox = slide.shapes.add_textbox(
            left + Inches(0.12), Inches(6.0), Inches(3.8), Inches(1.1))
        tf = txBox.text_frame
        tf.word_wrap = True
        for j, b in enumerate(bullets):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = b
            run.font.name = "Trebuchet MS"
            run.font.size = Pt(9)
            run.font.color.rgb = DARK_GRAY

    return slide


def slide_scheduling_rules(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, "Scheduling Rules & Constraints",
               "Logic applied by the engine when building each patient's timeline")
    footer_bar(slide, page, total)

    rules = [
        (BMS_PURPLE, "Priority-Based Ordering",
         ["Patients are sorted by priority before scheduling begins.",
          "Lower priority number = scheduled first.",
          "Affects which patient gets the 'next free' ROMAG slot."]),
        (BLUE, "Joint Incubator + ROMAG Assignment",
         ["Patient introduction requires BOTH a free incubator AND a free",
          "ROMAG + Romag Robot simultaneously.",
          "Engine tries all candidate incubators, picks the combination that",
          "gives the earliest UP_TSA start — minimising patient queue time."]),
        (TEAL, "INC Biological Minimum (2 h)",
         ["Each incubation step has a hard 2-hour minimum.",
          "Patient cannot advance to the next step before 2 h has elapsed,",
          "even if equipment is available earlier.",
          "Actual INC time extends retroactively to the next step's start."]),
        (ORANGE, "DHWF → FILL Chaining (No Hold)",
         ["Once material reaches the Harvest & Formulation (UP_DHWF) step,",
          "it must flow DIRECTLY into Fill (UP_FILL) — no intermediate hold.",
          "Scheduler iteratively aligns DHWF start so DHWF_end == FILL_start.",
          "If Fill is not available, DHWF is pushed later rather than held."]),
        (GREEN, "Continuous Incubator Reservation",
         ["The incubator assigned to a patient is reserved from UP_TSA start",
          "all the way to UP_DHWF start — even during non-INC steps.",
          "This reflects physical reality: the incubator cannot be reused",
          "mid-process for another patient."]),
        (RED, "47-Patient Concurrent Cap",
         ["Only 47 patients can be active simultaneously (one incubator each).",
          "Patient 48+ is queued; introduced when an incubator and Romag Robot",
          "are both simultaneously free.",
          "Weekly intake ≈ 47 incubators ÷ avg occupancy days × 7."]),
    ]

    for i, (color, title, bullets) in enumerate(rules):
        row, col = divmod(i, 2)
        left = Inches(0.35) + col * Inches(6.5)
        top  = Inches(1.35) + row * Inches(1.92)

        add_rect(slide, left, top, Inches(6.2), Inches(1.82),
                 LIGHT_GRAY, border_color=color, border_pt=2)
        # Colour strip
        add_rect(slide, left, top, Inches(0.18), Inches(1.82), color)
        add_text_box(slide, title,
                     left + Inches(0.28), top + Inches(0.08),
                     Inches(5.8), Inches(0.38),
                     font_size=12, bold=True, color=color)
        txBox = slide.shapes.add_textbox(
            left + Inches(0.28), top + Inches(0.5),
            Inches(5.8), Inches(1.2))
        tf = txBox.text_frame
        tf.word_wrap = True
        for j, b in enumerate(bullets):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = ("• " if b else "") + b
            run.font.name = "Trebuchet MS"
            run.font.size = Pt(9.5)
            run.font.color.rgb = DARK_GRAY

    return slide


def slide_incubator(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, "Incubator Pool Management",
               "47 incubators, one per patient, reserved for the full upstream window")
    footer_bar(slide, page, total)

    # Timeline diagram
    add_text_box(slide, "Patient Incubator Timeline (one patient)",
                 Inches(0.35), Inches(1.3), Inches(12.5), Inches(0.35),
                 font_size=12, bold=True, color=BMS_DARK)

    tl_top   = Inches(1.75)
    tl_left  = Inches(0.5)
    tl_width = Inches(12.3)

    # Background track
    add_rect(slide, tl_left, tl_top + Inches(0.32), tl_width, Inches(0.22), LIGHT_GRAY)

    # Phases along the timeline (proportional to ~160 h total)
    total_h = 160
    phases = [
        ("UP_TSA",   4,   BLUE,       False),
        ("INC 1",   48,   TEAL,        True),
        ("BSD",      2,   ORANGE,     False),
        ("INC 2",   24,   TEAL,        True),
        ("UP_TxD",   6,   BMS_PURPLE, False),
        ("INC 3",   48,   TEAL,        True),
        ("BSD",      2,   ORANGE,     False),
        ("INC 4",   24,   TEAL,        True),
    ]

    cursor = 0
    for code, hrs, color, is_inc in phases:
        w = tl_width * hrs / total_h
        x = tl_left + tl_width * cursor / total_h
        add_rect(slide, x, tl_top + Inches(0.32), w, Inches(0.22), color)
        if w > Inches(0.5):
            add_text_box(slide, code, x, tl_top + Inches(0.36), w, Inches(0.18),
                         font_size=7, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        cursor += hrs

    # Incubator reservation bar
    add_rect(slide, tl_left, tl_top, tl_width * cursor / total_h,
             Inches(0.28), BMS_PURPLE)
    add_text_box(slide, "← Incubator Reserved (UP_TSA start → UP_DHWF start)",
                 tl_left + Inches(0.1), tl_top + Inches(0.04),
                 tl_width - Inches(0.2), Inches(0.22),
                 font_size=8.5, bold=True, color=WHITE)

    # DHWF + FILL (after incubator released)
    x_dhwf = tl_left + tl_width * cursor / total_h
    for code, hrs, color in [("UP_DHWF", 8, GREEN), ("UP_FILL", 3, RED)]:
        w = tl_width * hrs / total_h
        add_rect(slide, x_dhwf, tl_top + Inches(0.32), w, Inches(0.22), color)
        add_text_box(slide, code, x_dhwf, tl_top + Inches(0.36), w, Inches(0.18),
                     font_size=7, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        x_dhwf += w

    # Time axis labels
    for h in [0, 40, 80, 120, 160]:
        x = tl_left + tl_width * h / total_h
        add_rect(slide, x, tl_top + Inches(0.56), Inches(0.01), Inches(0.1), MID_GRAY)
        add_text_box(slide, f"{h}h", x - Inches(0.15), tl_top + Inches(0.66),
                     Inches(0.3), Inches(0.18), font_size=7.5, color=MID_GRAY,
                     align=PP_ALIGN.CENTER)

    # Key insight boxes
    insights = [
        (BMS_PURPLE, "Reservation Span",
         "~7–10 days per patient\n(UP_TSA start → UP_DHWF start)\nIncubator cannot be reused during this window"),
        (TEAL, "Biological Minimums",
         "Each INC step: ≥ 2 h minimum\nPatient stays longer if next step's\nequipment is not yet free"),
        (ORANGE, "Queue Trigger",
         "Patient N+1 may enter the pool\nonly when Incubator AND ROMAG+\nRobot are simultaneously free"),
        (BLUE, "Weekly Intake Formula",
         "Patients/week =\n  47 incubators ÷ avg_days × 7\nReflects true slot turnover rate"),
    ]

    for i, (color, title, text) in enumerate(insights):
        left = Inches(0.35) + i * Inches(3.22)
        top  = Inches(3.05)
        add_rect(slide, left, top, Inches(3.05), Inches(3.9),
                 LIGHT_GRAY, border_color=color, border_pt=2)
        add_rect(slide, left, top, Inches(3.05), Inches(0.38), color)
        add_text_box(slide, title, left + Inches(0.1), top + Inches(0.05),
                     Inches(2.85), Inches(0.3),
                     font_size=11, bold=True, color=WHITE)
        add_text_box(slide, text, left + Inches(0.1), top + Inches(0.5),
                     Inches(2.85), Inches(3.3),
                     font_size=10.5, color=DARK_GRAY)

    return slide


def slide_efficiency(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, "Equipment Efficiency Factor",
               "OEE-style multiplier that scales scheduled durations and throughput projections")
    footer_bar(slide, page, total)

    # Concept explanation
    add_rect(slide, Inches(0.35), Inches(1.3), Inches(12.65), Inches(1.35),
             rgb(0xF0, 0xE8, 0xF9), border_color=BMS_PURPLE, border_pt=1.5)
    add_text_box(slide, "What is Equipment Efficiency?",
                 Inches(0.5), Inches(1.38), Inches(8), Inches(0.38),
                 font_size=13, bold=True, color=BMS_PURPLE)
    add_text_box(slide,
                 "Equipment efficiency (1–100%) represents how productively a piece of equipment "
                 "operates in practice — accounting for planned downtime, cleaning cycles, "
                 "maintenance, changeovers, and minor stoppages (analogous to OEE: "
                 "Overall Equipment Effectiveness).\n"
                 "A 100% efficient machine runs continuously at theoretical capacity. "
                 "An 80% efficient machine delivers only 80% of that capacity.",
                 Inches(0.5), Inches(1.78), Inches(12.3), Inches(0.8),
                 font_size=10.5, color=DARK_GRAY)

    # Formula box
    add_rect(slide, Inches(0.35), Inches(2.8), Inches(6.0), Inches(1.1),
             BMS_DARK, border_color=BMS_DARK, border_pt=0)
    add_text_box(slide, "Duration Scaling Formula",
                 Inches(0.5), Inches(2.85), Inches(5.7), Inches(0.32),
                 font_size=11, bold=True, color=BMS_PURPLE)
    add_text_box(slide,
                 "Effective Duration = ⌈ Nominal Duration ÷ (Efficiency / 100) ⌉",
                 Inches(0.5), Inches(3.18), Inches(5.7), Inches(0.55),
                 font_size=12, bold=True, color=WHITE)

    # Example table
    add_text_box(slide, "Example — UP_TSA (ROMAG, nominal 240 min):",
                 Inches(6.55), Inches(2.82), Inches(6.5), Inches(0.35),
                 font_size=11, bold=True, color=BMS_DARK)
    ex_data = [
        ("Efficiency", "Scheduled Duration", "Extra Time"),
        ("100%",       "240 min  (4.0 h)",   "+0 min"),
        ("90%",        "267 min  (4.5 h)",   "+27 min"),
        ("80%",        "300 min  (5.0 h)",   "+60 min"),
        ("70%",        "343 min  (5.7 h)",   "+103 min"),
        ("60%",        "400 min  (6.7 h)",   "+160 min"),
    ]
    for r, row in enumerate(ex_data):
        bg = BMS_MID if r == 0 else (LIGHT_GRAY if r % 2 == 0 else WHITE)
        tc = WHITE if r == 0 else DARK_GRAY
        bold = r == 0
        col_ws = [Inches(1.5), Inches(2.2), Inches(1.8)]
        col_xs = [Inches(6.55), Inches(8.1), Inches(10.35)]
        for val, cw, cx in zip(row, col_ws, col_xs):
            add_rect(slide, cx, Inches(3.2) + r * Inches(0.45), cw, Inches(0.45),
                     bg, border_color=MID_GRAY, border_pt=0.4)
            add_text_box(slide, val, cx + Inches(0.08),
                         Inches(3.22) + r * Inches(0.45),
                         cw - Inches(0.1), Inches(0.38),
                         font_size=10, bold=bold, color=tc)

    # Effect chain
    add_text_box(slide, "How Efficiency Flows Through the System",
                 Inches(0.35), Inches(4.08), Inches(12.5), Inches(0.38),
                 font_size=13, bold=True, color=BMS_DARK)

    chain = [
        (BMS_PURPLE, "User sets\nEfficiency %\nper equipment\ntype in Master\nData"),
        (BLUE,       "Engine inflates\nstep duration:\nnominal ÷ eff.\nAll steps on\nthat type"),
        (TEAL,       "Schedule\nblocks are\nlonger → later\nstarts for\nsubsequent steps"),
        (ORANGE,     "Target_duration_\nmin on each\nbatch reflects\nefficiency-adjusted\nduration"),
        (GREEN,      "Theoretical\nmax & AMS\nthroughput are\nautomatically\nreduced"),
        (RED,        "Bottleneck\nidentification\nshifts to reflect\ntrue constrained\ncapacity"),
    ]
    for i, (color, text) in enumerate(chain):
        left = Inches(0.35) + i * Inches(2.15)
        add_rect(slide, left, Inches(4.55), Inches(2.0), Inches(2.25),
                 LIGHT_GRAY, border_color=color, border_pt=2)
        add_rect(slide, left, Inches(4.55), Inches(2.0), Inches(0.12), color)
        add_text_box(slide, text, left + Inches(0.1), Inches(4.72),
                     Inches(1.8), Inches(2.0),
                     font_size=9.5, color=DARK_GRAY)
        if i < len(chain) - 1:
            add_text_box(slide, "→",
                         left + Inches(2.0), Inches(5.5),
                         Inches(0.2), Inches(0.3),
                         font_size=14, bold=True, color=MID_GRAY,
                         align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(0.35), Inches(6.92), Inches(12.65), Inches(0.1),
             BMS_PURPLE)
    add_text_box(slide,
                 "Note: The biological minimum (2 h INC floor) is NOT efficiency-adjusted — "
                 "it is a process constraint, not an equipment characteristic.",
                 Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.28),
                 font_size=9, italic=True, color=DARK_GRAY)

    return slide


def slide_throughput(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, "Throughput & Capacity Metrics",
               "How SIRIUS estimates annual production capacity from the live schedule")
    footer_bar(slide, page, total)

    metrics = [
        (BMS_PURPLE, "Avg Cycle Time",
         "Average wall-clock time from the first step (UP_TSA start) to "
         "the last step (UP_FILL end) across all scheduled patients.\n\n"
         "Unit: days\n"
         "Reflects real scheduling gaps, not just sum of step durations."),
        (BLUE, "Patients / Week\n(Weekly Intake)",
         "Maximum rate at which new patients can be introduced, given the "
         "incubator pool constraint.\n\n"
         "Formula:\n  47 incubators ÷ avg_occupancy_days × 7\n\n"
         "Occupancy = UP_TSA start → UP_DHWF start per patient."),
        (TEAL, "AMS Throughput\n(Annual Estimate)",
         "Annual capacity assuming AMS keeps the patient queue full at all times.\n\n"
         "Formula:\n  theoretical_max × bottleneck_util%\n\n"
         "Uses observed utilisation patterns to avoid distortion "
         "from small sample windows."),
        (GREEN, "Theoretical Max",
         "Absolute ceiling if the bottleneck equipment ran at 100% utilisation, 24/7.\n\n"
         "Formula:\n  n_units × 525,600 min/yr ÷ min_per_patient\n\n"
         "Identifies which equipment type is the binding constraint."),
    ]

    for i, (color, title, text) in enumerate(metrics):
        left = Inches(0.35) + i * Inches(3.22)
        top  = Inches(1.35)
        add_rect(slide, left, top, Inches(3.05), Inches(4.0),
                 LIGHT_GRAY, border_color=color, border_pt=2)
        add_rect(slide, left, top, Inches(3.05), Inches(0.42), color)
        add_text_box(slide, title, left + Inches(0.1), top + Inches(0.05),
                     Inches(2.85), Inches(0.38),
                     font_size=11, bold=True, color=WHITE)
        add_text_box(slide, text, left + Inches(0.1), top + Inches(0.52),
                     Inches(2.85), Inches(3.35),
                     font_size=9.5, color=DARK_GRAY)

    # Bottleneck explanation
    add_text_box(slide, "Bottleneck Identification",
                 Inches(0.35), Inches(5.5), Inches(12.5), Inches(0.38),
                 font_size=13, bold=True, color=BMS_DARK)

    bn_text = (
        "For each equipment type, SIRIUS computes: Theoretical_PPY = N_units × 525,600 ÷ mins_per_patient.\n"
        "The equipment type with the LOWEST theoretical_PPY is the bottleneck — "
        "it limits the entire system regardless of how much capacity other types have.\n"
        "When equipment efficiency is reduced, that type's theoretical_PPY falls further, "
        "potentially changing the bottleneck identity."
    )
    add_rect(slide, Inches(0.35), Inches(5.95), Inches(12.65), Inches(1.0),
             rgb(0xE8, 0xF4, 0xFF), border_color=BLUE, border_pt=1)
    add_text_box(slide, bn_text,
                 Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.95),
                 font_size=10, color=DARK_GRAY)

    return slide


def slide_bottleneck(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, "Bottleneck Detection & Robot Efficiency",
               "Identifying whether robots are constraints and quantifying their scheduling impact")
    footer_bar(slide, page, total)

    # Left: metrics explained
    add_text_box(slide, "Robot Efficiency Metrics",
                 Inches(0.35), Inches(1.3), Inches(6.0), Inches(0.38),
                 font_size=13, bold=True, color=BMS_DARK)

    metric_items = [
        (BMS_PURPLE, "Utilisation %",
         "Fraction of the total schedule span during which\nthe robot is actively operating.\n"
         "High utilisation → robot is working a lot.\n"
         "High utilisation alone does NOT mean it is a bottleneck."),
        (BLUE, "Idle %",
         "Complement of utilisation. Robot is powered on\nbut not performing any operation.\n"
         "Helps identify if shift patterns or gaps\nare leaving capacity on the table."),
        (ORANGE, "Contention Events",
         "Number of steps where the patient (or next patient)\nwas READY but the robot was BUSY.\n"
         "Includes both within-patient waits and\npatient-introduction delays (first step)."),
        (RED, "Patient Intros Delayed",
         "Subset of contention events where the delayed\nstep was step 0 (UP_TSA).\n"
         "Directly measures how many patients waited\nto start because the Romag Robot was occupied."),
        (TEAL, "Total & Avg Wait",
         "Total minutes of schedule delay attributable\nto robot contention.\n"
         "Avg = total ÷ events — shows severity\nof each individual delay."),
        (GREEN, "Scheduling Influence",
         "Qualitative label:\n"
         "• Not a constraint\n"
         "• Minor influence  (1–3 events)\n"
         "• Moderate bottleneck  (4–8 events)\n"
         "• Major bottleneck  (>8 events)\n"
         "• Bottleneck for patient introduction  (≥50% intros delayed)"),
    ]

    for i, (color, title, text) in enumerate(metric_items):
        row, col = divmod(i, 2)
        left = Inches(0.35) + col * Inches(3.1)
        top  = Inches(1.78) + row * Inches(1.72)
        add_rect(slide, left, top, Inches(2.95), Inches(1.62),
                 LIGHT_GRAY, border_color=color, border_pt=1.5)
        add_text_box(slide, title, left + Inches(0.1), top + Inches(0.07),
                     Inches(2.75), Inches(0.32),
                     font_size=10, bold=True, color=color)
        add_text_box(slide, text, left + Inches(0.1), top + Inches(0.42),
                     Inches(2.75), Inches(1.1),
                     font_size=8.5, color=DARK_GRAY)

    # Right: detection logic
    add_rect(slide, Inches(6.65), Inches(1.3), Inches(6.3), Inches(5.85),
             BMS_DARK, border_color=BMS_DARK, border_pt=0)
    add_text_box(slide, "How Contention is Detected",
                 Inches(6.8), Inches(1.38), Inches(6.0), Inches(0.38),
                 font_size=13, bold=True, color=BMS_PURPLE)

    detection_steps = [
        ("1", BMS_PURPLE,
         "For every batch that uses a robot, identify when the patient was 'ready'.\n"
         "For the first step (UP_TSA): ready = global schedule start.\n"
         "For subsequent steps: ready = previous step's actual end."),
        ("2", BLUE,
         "Compare the batch's actual start to the ready time.\n"
         "If actual_start > ready_time, a gap exists — the patient was waiting."),
        ("3", TEAL,
         "Check whether the robot was BUSY during the gap.\n"
         "If yes → robot caused the delay → contention event.\n"
         "Count += 1, accumulate wait minutes."),
        ("4", ORANGE,
         "Flag separately if step_index == 0 → patient introduction delay.\n"
         "If ≥50% of patients had their introduction delayed, label:\n"
         "'Bottleneck for patient introduction'."),
    ]

    for i, (num, color, text) in enumerate(detection_steps):
        top = Inches(1.88) + i * Inches(1.3)
        add_rect(slide, Inches(6.8), top, Inches(0.38), Inches(0.38), color)
        add_text_box(slide, num, Inches(6.8), top,
                     Inches(0.38), Inches(0.38),
                     font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, text,
                     Inches(7.28), top, Inches(5.5), Inches(1.2),
                     font_size=9.5, color=WHITE)

    return slide


def slide_controls(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, "Controls & Monitoring",
               "Simulation engine, live tracking, and system integrations")
    footer_bar(slide, page, total)

    # Top row: core controls
    controls = [
        (BMS_PURPLE, "Schedule Generation",
         ["Priority-sorted patients",
          "Joint incubator + ROMAG optimisation",
          "Efficiency-adjusted durations",
          "Chained DHWF → FILL constraint",
          "Retroactive INC extension",
          "Auto-assigns order numbers"]),
        (BLUE, "Simulation Engine",
         ["Mark batches as IN_PROGRESS / COMPLETE",
          "Auto-detect delays vs max duration threshold",
          "Trigger reschedule when delay > threshold",
          "Early completion → pull forward downstream",
          "Auto-complete segments proportionally on batch end",
          "Robot calendars updated in real time"]),
        (TEAL, "Live Mode",
         ["Clock-driven automatic status updates",
          "Auto-start batches after grace period",
          "Auto-reschedule on detected delays",
          "Runs on every Streamlit page render",
          "Toggle on/off without losing schedule state"]),
        (ORANGE, "Session Persistence",
         ["Full schedule saved to JSON on every change",
          "Patients, batches, sequences, equipment",
          "Equipment efficiency settings persisted",
          "Reload on browser refresh — no data loss",
          "Export / import session files"]),
    ]

    for i, (color, title, bullets) in enumerate(controls):
        left = Inches(0.35) + i * Inches(3.22)
        top  = Inches(1.35)
        add_rect(slide, left, top, Inches(3.05), Inches(3.5),
                 LIGHT_GRAY, border_color=color, border_pt=2)
        add_rect(slide, left, top, Inches(3.05), Inches(0.38), color)
        add_text_box(slide, title, left + Inches(0.1), top + Inches(0.05),
                     Inches(2.85), Inches(0.32),
                     font_size=11, bold=True, color=WHITE)
        txBox = slide.shapes.add_textbox(
            left + Inches(0.1), top + Inches(0.5),
            Inches(2.85), Inches(2.85))
        tf = txBox.text_frame
        tf.word_wrap = True
        for j, b in enumerate(bullets):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = "• " + b
            run.font.name = "Trebuchet MS"
            run.font.size = Pt(9.5)
            run.font.color.rgb = DARK_GRAY

    # Bottom: integrations
    add_text_box(slide, "System Integrations",
                 Inches(0.35), Inches(5.0), Inches(12.5), Inches(0.38),
                 font_size=13, bold=True, color=BMS_DARK)

    integrations = [
        (GREEN, "SAP Interface",
         "Reads patient orders (batch IDs, priorities, order numbers) from SAP. "
         "Status changes can be written back on batch completion."),
        (RED, "DeltaV MES Interface",
         "Monitors batch execution against the schedule. "
         "Provides actual start/end timestamps for live deviation detection."),
        (BMS_PURPLE, "Gantt Visualisation",
         "Interactive Plotly Gantt chart with datetime slider for time range zoom. "
         "Equipment view and Patient view. Colour-coded by phase."),
        (BLUE, "Master Data Editor",
         "In-app editing of process sequences, step durations, robot types, "
         "sigma thresholds, and equipment efficiency — no code changes needed."),
    ]

    for i, (color, title, text) in enumerate(integrations):
        left = Inches(0.35) + i * Inches(3.22)
        top  = Inches(5.48)
        add_rect(slide, left, top, Inches(3.05), Inches(1.5),
                 LIGHT_GRAY, border_color=color, border_pt=1.5)
        add_text_box(slide, title, left + Inches(0.1), top + Inches(0.06),
                     Inches(2.85), Inches(0.32),
                     font_size=10, bold=True, color=color)
        add_text_box(slide, text, left + Inches(0.1), top + Inches(0.42),
                     Inches(2.85), Inches(1.0),
                     font_size=9, color=DARK_GRAY)

    return slide


def slide_summary(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, BMS_DARK)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), BMS_PURPLE)
    add_rect(slide, 0, Inches(7.42), SLIDE_W, Inches(0.08), BMS_PURPLE)

    add_text_box(slide, "Key Takeaways",
                 Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
                 font_size=28, bold=True, color=WHITE)

    takeaways = [
        (BMS_PURPLE, "Patient-Centric, End-to-End",
         "Every patient batch follows a strict 10-step sequence from T-cell activation "
         "to fill & finish. No pooling. Full traceability from incubator assignment "
         "to robot calendar."),
        (BLUE, "Joint Resource Optimisation",
         "Patient introduction is gated by simultaneous availability of an incubator, "
         "a ROMAG, and the Romag Robot. The scheduler co-optimises all three to "
         "minimise queue time and maximise throughput."),
        (TEAL, "47-Incubator Pool is the Primary Capacity Lever",
         "The incubator count determines maximum concurrent patients. Occupancy per "
         "patient spans the full upstream window (~7–10 days). Increasing the pool "
         "or reducing cycle time directly improves weekly intake."),
        (ORANGE, "Robot Utilisation ≠ Robot Constraint",
         "A robot can be 60% utilised without causing any delays, or 30% utilised "
         "and still be the reason every new patient waits. Contention events and "
         "patient-introduction delays are the real scheduling-impact metrics."),
        (GREEN, "Efficiency Factor Cascades",
         "Lowering equipment efficiency inflates step durations, reduces theoretical "
         "throughput, and may shift the bottleneck identity — all automatically "
         "reflected in the schedule and capacity metrics."),
        (RED, "Controls Close the Loop",
         "Simulation mode, live clock tracking, and auto-reschedule keep the schedule "
         "aligned with reality. SAP and DeltaV integrations pull and push actual "
         "execution data without manual intervention."),
    ]

    for i, (color, title, text) in enumerate(takeaways):
        row, col = divmod(i, 2)
        left = Inches(0.4) + col * Inches(6.45)
        top  = Inches(1.1) + row * Inches(1.95)
        add_rect(slide, left, top, Inches(6.1), Inches(1.8),
                 rgb(0x23, 0x2B, 0x3B), border_color=color, border_pt=2)
        add_rect(slide, left, top, Inches(0.16), Inches(1.8), color)
        add_text_box(slide, title,
                     left + Inches(0.28), top + Inches(0.1),
                     Inches(5.7), Inches(0.38),
                     font_size=11.5, bold=True, color=color)
        add_text_box(slide, text,
                     left + Inches(0.28), top + Inches(0.52),
                     Inches(5.7), Inches(1.15),
                     font_size=9.5, color=WHITE)

    add_text_box(slide, "Questions?  Contact the Cell Therapy Operations Scheduling Team.",
                 Inches(0.6), Inches(7.08), Inches(12), Inches(0.32),
                 font_size=10, color=MID_GRAY, align=PP_ALIGN.CENTER)

    return slide


# ── Main ───────────────────────────────────────────────────────────────────────

def build_presentation():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    TOTAL = 10  # title + agenda + 8 content slides

    slide_title(prs)
    slide_agenda(prs, 2, TOTAL)
    slide_cart_overview(prs, 3, TOTAL)
    slide_process(prs, 4, TOTAL)
    slide_equipment(prs, 5, TOTAL)
    slide_scheduling_rules(prs, 6, TOTAL)
    slide_incubator(prs, 7, TOTAL)
    slide_efficiency(prs, 8, TOTAL)
    slide_throughput(prs, 9, TOTAL)
    slide_bottleneck(prs, 10, TOTAL)
    # slide_controls goes to 11 but let's include it
    # re-total
    return prs


def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slides_fns = [
        slide_title,
        slide_agenda,
        slide_cart_overview,
        slide_process,
        slide_equipment,
        slide_scheduling_rules,
        slide_incubator,
        slide_efficiency,
        slide_throughput,
        slide_bottleneck,
        slide_controls,
        slide_summary,
    ]
    total = len(slides_fns)

    for i, fn in enumerate(slides_fns):
        page = i + 1
        if fn == slide_title:
            fn(prs)
        else:
            fn(prs, page, total)

    out = "SIRIUS_Scheduler_Overview.pptx"
    prs.save(out)
    print(f"Saved: {out}  ({total} slides)")


if __name__ == "__main__":
    main()
